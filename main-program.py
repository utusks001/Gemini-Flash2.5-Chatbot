import os
from io import BytesIO
import requests
import streamlit as st
from dotenv import load_dotenv

# File parsing
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation

# Data analysis
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns

# LangChain
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

# -------------------------
# Config
# -------------------------
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")
OCR_SPACE_API_KEY = os.getenv("OCR_SPACE_API_KEY")

st.set_page_config(page_title="Chatbot + Excel Analysis", page_icon="🤖", layout="wide")

# -------------------------
# Session state
# -------------------------
if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "indexed_files" not in st.session_state:
    st.session_state.indexed_files = []
if "dataframes" not in st.session_state:
    st.session_state.dataframes = {}
if "last_uploaded" not in st.session_state:
    st.session_state.last_uploaded = []

# -------------------------
# Helpers for DataFrame
# -------------------------
def safe_describe(df):
    try:
        return df.describe(include="all", datetime_is_numeric=True)
    except TypeError:
        return df.describe(include="all")

def df_info_text(df):
    import io
    buf = io.StringIO()
    df.info(buf=buf)
    return buf.getvalue()

def df_to_index_text(df, filename, sheet_name):
    rows, cols = df.shape
    stats = safe_describe(df).transpose().reset_index().to_string(index=False)
    sample = df.head(20).to_csv(index=False)
    return f"DATAFRAME — file={filename}, sheet={sheet_name}\nshape: {rows}x{cols}\n{stats}\nSAMPLE:\n{sample}"

# -------------------------
# Extractors
# -------------------------
def extract_text_from_pdf(file_bytes):
    text = ""
    try:
        reader = PdfReader(file_bytes)
        for page in reader.pages:
            if page.extract_text():
                text += page.extract_text() + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak PDF: {e}")
    return text

def extract_text_from_txt(file_bytes):
    try:
        return file_bytes.read().decode("utf-8", errors="ignore")
    except Exception as e:
        st.warning(f"⚠️ Gagal baca TXT: {e}")
        return ""

def extract_text_from_docx(file_bytes):
    text = ""
    try:
        doc = DocxDocument(file_bytes)
        for p in doc.paragraphs:
            if p.text:
                text += p.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak DOCX: {e}")
    return text

def extract_text_from_pptx(file_bytes):
    text = ""
    try:
        prs = PptxPresentation(file_bytes)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak PPTX: {e}")
    return text

def extract_text_from_image(file_bytes, filename="upload.png"):
    if not OCR_SPACE_API_KEY:
        st.warning("⚠️ OCR_SPACE_API_KEY tidak ditemukan di .env — OCR image dinonaktifkan.")
        return ""
    try:
        file_bytes.seek(0)
        data = file_bytes.read()
        if not data:
            st.warning(f"⚠️ Gambar {filename} kosong atau gagal terbaca.")
            return ""

        resp = requests.post(
            "https://api.ocr.space/parse/image",
            files={"file": (filename, BytesIO(data), "image/png")},
            data={
                "apikey": OCR_SPACE_API_KEY,
                "language": "eng",
                "isOverlayRequired": False
            },
            timeout=60
        )

        # ✅ Debug patch
        try:
            result = resp.json()
        except Exception:
            st.error("⚠️ OCR.Space tidak balas JSON.\n\nRespons mentah:\n\n" + resp.text[:500])
            return ""

        if not isinstance(result, dict):
            st.error("⚠️ OCR.Space respons bukan JSON valid.\n\nRespons mentah:\n\n" + str(result)[:500])
            return ""

        if result.get("IsErroredOnProcessing"):
            st.error("⚠️ OCR.Space error: " + str(result.get("ErrorMessage", ['Unknown error'])))
            return ""

        parsed = [p.get("ParsedText", "") for p in result.get("ParsedResults", []) if isinstance(p, dict)]
        return "\n".join(parsed).strip()
    except Exception as e:
        st.error(f"⚠️ OCR error: {e}")
        return ""

def extract_text_from_csv(file_bytes, filename):
    try:
        df = pd.read_csv(file_bytes)
        st.session_state.dataframes[filename] = {"sheets": {"CSV": df}}
        return df_to_index_text(df, filename, "CSV")
    except Exception as e:
        st.warning(f"⚠️ Gagal baca CSV {filename}: {e}")
        return ""

def extract_text_from_excel(file_bytes, filename):
    text_parts = []
    try:
        xls = pd.ExcelFile(file_bytes)
        sheet_map = {}
        for s in xls.sheet_names:
            try:
                df = xls.parse(s)
                sheet_map[s] = df
                text_parts.append(df_to_index_text(df, filename, s))
            except Exception as se:
                st.warning(f"⚠️ Gagal parse sheet '{s}' di {filename}: {se}")
        if sheet_map:
            st.session_state.dataframes[filename] = {"sheets": sheet_map}
        return "\n\n---\n\n".join(text_parts)
    except Exception as e:
        st.warning(f"⚠️ Gagal baca Excel {filename}: {e}")
        return ""

# -------------------------
# Dispatcher
# -------------------------
def extract_text_from_file(uploaded_file):
    name = uploaded_file.name
    lname = name.lower()
    raw = uploaded_file.getvalue()

    if lname.endswith(".pdf"):
        return extract_text_from_pdf(BytesIO(raw))
    if lname.endswith(".txt"):
        return extract_text_from_txt(BytesIO(raw))
    if lname.endswith(".docx"):
        return extract_text_from_docx(BytesIO(raw))
    if lname.endswith(".pptx"):
        return extract_text_from_pptx(BytesIO(raw))
    if lname.endswith(".csv"):
        return extract_text_from_csv(BytesIO(raw), name)
    if lname.endswith((".xlsx", ".xls")):
        return extract_text_from_excel(BytesIO(raw), name)
    if lname.endswith((".jpg", ".jpeg", ".png", ".gif", ".bmp", ".jfif")):
        return extract_text_from_image(BytesIO(raw), filename=name)

    st.warning(f"⚠️ Format file `{name}` tidak didukung.")
    return ""

# -------------------------
# Build docs & FAISS
# -------------------------
EMBEDDINGS = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
SPLITTER = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=120)

def build_documents_from_uploads(files):
    docs = []
    for f in files:
        text = extract_text_from_file(f)
        if text.strip():
            chunks = SPLITTER.split_text(text)
            for i, chunk in enumerate(chunks):
                docs.append(Document(page_content=chunk, metadata={"source": f.name, "chunk_id": i}))
    return docs

def build_faiss_from_documents(docs):
    if not docs:
        return None
    return FAISS.from_documents(docs, embedding=EMBEDDINGS)

# -------------------------
# Auto-analysis
# -------------------------
def auto_analyze_dataframe(df, filename, sheet_name):
    num_df = df.select_dtypes(include="number")

    st.markdown(f"### 📄 Analisa: {filename} — {sheet_name}")
    st.write("**Head (10):**")
    st.dataframe(df.head(10))
    st.write("**Tail (10):**")
    st.dataframe(df.tail(10))
    st.write("**describe():**")
    st.dataframe(safe_describe(df))
    st.write("**info():**")
    st.text(df_info_text(df))

    target_cols = [c for c in ["Sales", "Quantity", "Profit"] if c in df.columns]
    if target_cols:
        st.write("**Outlier Detection (Boxplot)**")
        for col in target_cols:
            fig, ax = plt.subplots(figsize=(5, 3))
            sns.boxplot(x=df[col].dropna(), ax=ax)
            ax.set_title(f"Outliers — {col}")
            st.pyplot(fig)

    if "Sales" in df.columns and "Profit" in df.columns:
        st.write("**Scatter Plot: Sales vs Profit**")
        fig, ax = plt.subplots(figsize=(5, 3))
        sns.scatterplot(x=df["Sales"], y=df["Profit"], ax=ax)
        st.pyplot(fig)

    if not num_df.empty:
        st.write("**Correlation Heatmap**")
        fig, ax = plt.subplots(figsize=(5, 3))
        sns.heatmap(num_df.corr(), annot=True, cmap="coolwarm", ax=ax)
        st.pyplot(fig)

# -------------------------
# UI
# -------------------------
st.title("🤖 Chatbot + Multi-files + Excel Analysis")

uploaded_files = st.sidebar.file_uploader(
    "Upload files",
    type=["pdf","txt","docx","pptx","jpg","jpeg","png","gif","bmp","jfif","csv","xls","xlsx"],
    accept_multiple_files=True
)

# ✅ Reset otomatis jika file baru diupload
if uploaded_files and uploaded_files != st.session_state.last_uploaded:
    st.session_state.vector_store = None
    st.session_state.indexed_files = []
    st.session_state.dataframes = {}
    st.session_state.last_uploaded = uploaded_files

    for f in uploaded_files:
        extract_text_from_file(f)

# ✅ Tombol reset manual
if st.sidebar.button("🧹 Reset Data"):
    st.session_state.vector_store = None
    st.session_state.indexed_files = []
    st.session_state.dataframes = {}
    st.session_state.last_uploaded = []
    st.sidebar.success("Data berhasil direset.")

if st.sidebar.button("🚀 Build Vector Store"):
    if uploaded_files:
        with st.spinner("Membangun vector store..."):
            docs = build_documents_from_uploads(uploaded_files)
            st.session_state.vector_store = build_faiss_from_documents(docs)
            st.session_state.indexed_files = [f.name for f in uploaded_files]
        st.sidebar.success("✅ Vector store terbangun")

if st.session_state.dataframes:
    st.subheader("📊 Analisa Excel/CSV")
    for fname, payload in st.session_state.dataframes.items():
        for sheet, df in payload["sheets"].items():
            auto_analyze_dataframe(df, fname, sheet)

# -------------------------
# Q&A Section
# -------------------------
st.subheader("💬 Tanya Jawab Dokumen")

model_choice = st.sidebar.radio("LLM Provider:", ["Gemini 2.5 Flash", "Groq Llama"])
prompt = st.text_input("Ajukan pertanyaan berdasarkan dokumen:")
ask_btn = st.button("Tanyakan")

def format_context(results):
    context_text = ""
    for i, d in enumerate(results):
        context_text += f"[{i+1}] {d.page_content[:500]}...\n\n"
    return context_text

if ask_btn:
    if not prompt.strip():
        st.warning("Masukkan pertanyaan terlebih dahulu.")
    elif st.session_state.vector_store is None:
        st.info("Belum ada vector store. Upload file dan klik 'Build Vector Store'.")
    else:
        with st.spinner("🔎 Mengambil konteks dari vector store..."):
            results = st.session_state.vector_store.similarity_search(prompt, k=5)

        context_text = format_context(results)
        system_instructions = (
            "Jawablah seakurat dan sedetil mungkin berdasarkan konteks berikut. "
            "Jika konteks berupa ringkasan tabel (CSV/Excel), gunakan metrik yang tersedia (shape, dtypes, missing, describe, sample). "
            "Jika jawaban tidak ada, katakan: "
            "\"Jawaban tidak tersedia dalam konteks yang diberikan, "
            "pelajari dan berikan tambahan informasi yang diperlukan sesuai konteks atau topik "
            "yang ditanyakan dari berbagai sumber atau referensi di internet dengan menyertakan link website.\" "
            "Sertakan referensi [angka] ke potongan konteks bila relevan."
        )

        composed_prompt = (
            f"{system_instructions}\n\n"
            f"=== KONTEX ===\n{context_text}\n\n"
            f"=== PERTANYAAN ===\n{prompt}\n\n"
            f"=== JAWABAN ==="
        )

        try:
            if model_choice.startswith("Gemini"):
                from langchain_google_genai import ChatGoogleGenerativeAI
                llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.2)
                resp = llm.invoke(composed_prompt)
            else:
                from langchain_groq import ChatGroq
                llm = ChatGroq(
                    model_name="llama-3.3-70b-versatile",
                    groq_api_key=GROQ_API_KEY,
                    temperature=0.2
                )
                resp = llm.invoke(composed_prompt)

            st.subheader("💬 Jawaban")
            st.write(getattr(resp, "content", str(resp)))

            # 🔎 Referensi
            if results:
                st.markdown("**📚 Referensi:**")
                for i, d in enumerate(results):
                    source = d.metadata.get("source", "Unknown Source")
                    st.markdown(f"[{i+1}] {source}")
        except Exception as e:
            st.error(f"❌ LLM error: {e}")
