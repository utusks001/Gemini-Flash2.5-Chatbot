# Streamlit UI
# program-utama.py
#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
Gemini / Groq Image‑enabled Multi‑File Chatbot
================================================
A Streamlit app that accepts PDFs, DOCX, PPTX, TXT **and** images, extracts text
( OCR for images ), builds a FAISS vector store, and answers questions using
Gemini or Groq’s Llama‑3.1‑405B‑Instruct.

Author:  OpenAI‑ChatGPT
License: MIT
"""

import os
import io
import base64
from typing import List, Optional

import streamlit as st
from dotenv import load_dotenv
from PIL import Image
import easyocr

# ---------------------------------
# LangChain imports
# ---------------------------------
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

# LLM wrappers
try:
    from langchain_google_genai import ChatGoogleGenerativeAI
except Exception:
    ChatGoogleGenerativeAI = None
try:
    from langchain_groq import ChatGroq
except Exception:
    ChatGroq = None

# ---------------------------------
# Environment / config
# ---------------------------------
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY")

if not (GOOGLE_API_KEY or GROQ_API_KEY):
    st.error(
        "❌ **No LLM credentials found**.\n"
        "Set either `GOOGLE_API_KEY` or `GROQ_API_KEY` in a `.env` file or Streamlit secrets."
    )
    st.stop()

# ---------------------------------
# Global objects
# ---------------------------------
EMBEDDINGS = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
SPLITTER = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=120)

# OCR reader – only load once
OCR_READER = easyocr.Reader(
    ["en"],  # you can add other languages if needed
    gpu=False,  # set True if you have a GPU and want faster OCR
    verbose=False,
)

# ---------------------------------
# File parsing helpers
# ---------------------------------
def extract_text_from_pdf(file_bytes: io.BytesIO) -> str:
    """Extract text from PDF using PyPDF2."""
    from PyPDF2 import PdfReader

    text = ""
    try:
        reader = PdfReader(file_bytes)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        st.warning(f"⚠️ PDF extraction error: {e}")
    return text


def extract_text_from_txt(file_bytes: io.BytesIO) -> str:
    """Extract text from a plain text file."""
    try:
        return file_bytes.read().decode("utf-8", errors="ignore")
    except Exception as e:
        st.warning(f"⚠️ TXT read error: {e}")
        return ""


def extract_text_from_docx(file_bytes: io.BytesIO) -> str:
    """Extract text from DOCX using python‑docx."""
    from docx import Document as DocxDocument

    text = ""
    try:
        file_bytes.seek(0)
        doc = DocxDocument(file_bytes)
        for p in doc.paragraphs:
            if p.text:
                text += p.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ DOCX extraction error: {e}")
    return text


def extract_text_from_pptx(file_bytes: io.BytesIO) -> str:
    """Extract text from PPTX using python‑pptx."""
    from pptx import Presentation as PptxPresentation

    text = ""
    try:
        file_bytes.seek(0)
        prs = PptxPresentation(file_bytes)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ PPTX extraction error: {e}")
    return text


def extract_text_from_image(file_bytes: io.BytesIO) -> str:
    """OCR a single image file using EasyOCR."""
    try:
        # read image with PIL
        image = Image.open(file_bytes).convert("RGB")
        # convert to numpy array
        import numpy as np

        img_np = np.array(image)
        # OCR
        results = OCR_READER.readtext(img_np, detail=0)  # detail=0 returns only text
        return "\n".join(results)
    except Exception as e:
        st.warning(f"⚠️ Image OCR error: {e}")
        return ""


def extract_text_from_file(uploaded_file) -> str:
    """Dispatch to the right extractor based on file extension."""
    name = uploaded_file.name.lower()
    raw_bytes = uploaded_file.read()
    bio = io.BytesIO(raw_bytes)

    if name.endswith(".pdf"):
        return extract_text_from_pdf(bio)
    elif name.endswith(".txt"):
        return extract_text_from_txt(bio)
    elif name.endswith(".docx"):
        return extract_text_from_docx(bio)
    elif name.endswith(".pptx"):
        return extract_text_from_pptx(bio)
    elif name.endswith((".bmp", ".png", ".jpg", ".jpeg", ".gif", ".jfif")):
        return extract_text_from_image(bio)
    else:
        st.warning(f"⚠️ Unsupported file type: `{uploaded_file.name}`")
        return ""


# ---------------------------------
# Build documents and vector store
# ---------------------------------
def build_documents_from_uploads(uploaded_files: List[st.UploadedFile]) -> List[Document]:
    """Return a list of LangChain Document objects for each chunk."""
    docs: List[Document] = []
    for f in uploaded_files:
        text = extract_text_from_file(f)
        if not text or not text.strip():
            continue
        chunks = SPLITTER.split_text(text)
        for i, chunk in enumerate(chunks):
            docs.append(
                Document(
                    page_content=chunk,
                    metadata={"source_file": f.name, "chunk_id": i},
                )
            )
    return docs


def build_faiss_from_documents(docs: List[Document]) -> Optional[FAISS]:
    """Create a FAISS vector store from the documents."""
    if not docs:
        return None
    return FAISS.from_documents(docs, embedding=EMBEDDINGS)


# ---------------------------------
# Prompt helpers
# ---------------------------------
def format_context(snippets: List[Document]) -> str:
    parts = []
    for idx, d in enumerate(snippets, start=1):
        src = d.metadata.get("source_file", "unknown")
        cid = d.metadata.get("chunk_id", "-")
        parts.append(
            f"[{idx}] ({src}#chunk-{cid})\n{d.page_content}"
        )
    return "\n\n---\n\n".join(parts)


def render_sources(snippets: List[Document]):
    with st.expander("🔎 Sumber konteks yang dipakai"):
        for i, d in enumerate(snippets, start=1):
            src = d.metadata.get("source_file", "unknown")
            cid = d.metadata.get("chunk_id", "-")
            preview = d.page_content[:300].replace("\n", " ")
            st.markdown(f"**[{i}]** **{src}** (chunk {cid})")
            st.caption(preview + ("..." if len(d.page_content) > 300 else ""))


# ---------------------------------
# LLM helper
# ---------------------------------
def get_llm():
    """Return a configured LLM wrapper based on available API key."""
    if GOOGLE_API_KEY and ChatGoogleGenerativeAI:
        return ChatGoogleGenerativeAI(
            model="gemini-2.5-flash",
            temperature=0.2,
            google_api_key=GOOGLE_API_KEY,
        )
    elif GROQ_API_KEY and ChatGroq:
        return ChatGroq(
            model="llama-3.1-405b-instruct",
            temperature=0.2,
            groq_api_key=GROQ_API_KEY,
        )
    else:
        st.error("❌ LLM library not available.")
        return None


# ---------------------------------
# Streamlit UI
# ---------------------------------
st.set_page_config(
    page_title="Gemini / Groq Image‑Chatbot – Multi‑File, OCR, FAISS",
    page_icon="🤖",
    layout="wide",
)

st.title("🤖 Gemini / Groq Image‑Chatbot – Multi‑File, OCR, FAISS")
st.markdown(
    """
Upload any number of files (PDF, TXT, DOCX, PPTX, **BMP, PNG, JPG, JPEG, GIF, JFIF**).  
The app will:
1. **Extract** text (OCR for images)  
2. **Chunk** the text  
3. **Embed** & **index** with FAISS  
4. **Answer** your questions using Gemini 2.5 Flash (Google) or Llama‑3.1‑405B‑Instruct (Groq)
"""
)

# Sidebar – file upload + build
st.sidebar.header("📂 Upload & Build")
uploaded_files = st.sidebar.file_uploader(
    "Upload files (pdf, txt, docx, pptx, bmp, png, jpg, jpeg, gif, jfif)",
    type=[
        "pdf",
        "txt",
        "docx",
        "pptx",
        "bmp",
        "png",
        "jpg",
        "jpeg",
        "gif",
        "jfif",
    ],
    accept_multiple_files=True,
)

build_btn = st.sidebar.button("🚀 Build Vector Store")
clear_btn = st.sidebar.button("🧹 Reset vector store")

if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "indexed_files" not in st.session_state:
    st.session_state.indexed_files = []

# Reset
if clear_btn:
    st.session_state.vector_store = None
    st.session_state.indexed_files = []
    st.sidebar.success("Vector store di‑reset.")

# Build
if build_btn:
    if not uploaded_files:
        st.sidebar.warning("Silakan upload minimal 1 file terlebih dahulu.")
    else:
        with st.spinner("📦 Memproses file & membuat vector store..."):
            docs = build_documents_from_uploads(uploaded_files)
            if not docs:
                st.sidebar.error(
                    "Tidak ada teks valid yang berhasil diekstrak. Pastikan file berisi teks atau gambar berisi teks."
                )
            else:
                vs = build_faiss_from_documents(docs)
                st.session_state.vector_store = vs
                st.session_state.indexed_files = [f.name for f in uploaded_files]
                st.sidebar.success(
                    f"Vector store terbangun. Dokumen terindeks: {len(st.session_state.indexed_files)} | Chunk total: {len(docs)}"
                )

# Show indexed files
if st.session_state.indexed_files:
    st.markdown("**Dokumen terindeks:**")
    st.write(" • " + "\n • ".join(st.session_state.indexed_files))

# Query area
st.subheader("🗣️ Tanya sesuatu...")
prompt = st.text_input(
    "Masukkan pertanyaan:",
    placeholder="Contoh: Ringkas dokumen tentang X…",
)
ask_btn = st.button("Tanyakan")

if ask_btn:
    if not prompt or not prompt.strip():
        st.warning("Masukkan pertanyaan terlebih dahulu.")
    elif st.session_state.vector_store is None:
        st.info(
            "Belum ada vector store. Upload file dan klik **Build Vector Store** di sidebar."
        )
    else:
        with st.spinner("🔎 Mencari konteks…"):
            results = st.session_state.vector_store.similarity_search(prompt, k=5)

        context_text = format_context(results)

        # System instructions – keep same as original
        system_instructions = (
            "Jawablah seakurat dan sedetil mungkin sehingga mudah dipahami berdasarkan konteks berikut. "
            "Jika jawaban tidak ada, katakan: "
            "\"Jawaban tidak tersedia dalam konteks yang diberikan, "
            "pelajari dan berikan tambahan informasi yang diperlukan sesuai konteks atau topik "
            "yang ditanyakan dari berbagai sumber atau referensi di internet dengan menyertakan link website.\" "
            "Sertakan referensi [angka] ke potongan konteks bila relevan."
        )

        composed_prompt = (
            f"{system_instructions}\n\n"
            f"=== KONTEX ===\n{context_text}\n\n"
            f"=== PERTANYAAN ===\n{prompt}\n\n"
            f"=== JAWABAN ==="
        )

        llm = get_llm()
        if llm is None:
            st.error("❌ LLM tidak dapat dipanggil.")
        else:
            try:
                with st.spinner("🤖 Menjawab…"):
                    response = llm.invoke(composed_prompt)

                st.subheader("💬 Jawaban")
                out_text = getattr(response, "content", None)
                if not out_text:
                    # fallback for newer langchain
                    out_text = str(response)
                st.write(out_text)

                render_sources(results)

            except Exception as e:
                st.error(f"❌ Error saat memanggil LLM: {e}")

# ---------------------------------
# Footer
# ---------------------------------
st.markdown("---")
st.markdown(
    """
**Disclaimer**  
*The OCR quality depends on the image resolution. For best results, upload clear, high‑resolution images.*  
*The answer is generated by a large language model and is not a substitute for professional advice.*
"""
)

