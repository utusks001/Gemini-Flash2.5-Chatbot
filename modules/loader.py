# Load & split dokumen + OCR image
# OCR.space + File Parser
# modules/loader.py

from io import BytesIO
from PIL import Image
import requests
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from pypdf import PdfReader
import streamlit as st
import os

splitter = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=120)
OCRSPACE_API_KEY = os.getenv("OCRSPACE_API_KEY", "helloworld")

def ocr_with_ocrspace(image_bytes):
    response = requests.post(
        'https://api.ocr.space/parse/image',
        files={'filename': image_bytes},
        data={'apikey': OCRSPACE_API_KEY, 'language': 'eng'}
    )
    result = response.json()
    return result['ParsedResults'][0]['ParsedText'] if 'ParsedResults' in result else ""

def extract_text_from_file(uploaded_file) -> str:
    name = uploaded_file.name.lower()
    raw = uploaded_file.read()
    bio = BytesIO(raw)

    try:
        if name.endswith(".pdf"):
            reader = PdfReader(bio)
            return "\n".join([p.extract_text() or "" for p in reader.pages])
        elif name.endswith(".txt"):
            return raw.decode("utf-8", errors="ignore")
        elif name.endswith(".docx"):
            doc = DocxDocument(bio)
            return "\n".join([p.text for p in doc.paragraphs])
        elif name.endswith(".pptx"):
            prs = PptxPresentation(bio)
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif name.endswith((".png", ".jpg", ".jpeg", ".bmp", ".gif", ".jfif")):
            return ocr_with_ocrspace(raw)
        else:
            st.warning(f"⚠️ Format tidak didukung: {uploaded_file.name}")
            return ""
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak {uploaded_file.name}: {e}")
        return ""

def build_documents_from_uploads(uploaded_files):
    docs = []
    for f in uploaded_files:
        text = extract_text_from_file(f)
        if not text.strip():
            continue
        chunks = splitter.split_text(text)
        for i, chunk in enumerate(chunks):
            docs.append(Document(page_content=chunk, metadata={"source_file": f.name, "chunk_id": i}))
    return docs
