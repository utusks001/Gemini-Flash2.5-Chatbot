# Multi-format Extractor
from io import BytesIO
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation

def extract_text_from_file(uploaded_file):
    name = uploaded_file.name.lower()
    raw = uploaded_file.read()
    bio = BytesIO(raw)
    text = ""

    if name.endswith(".pdf"):
        reader = PdfReader(bio)
        text = "\n".join([p.extract_text() for p in reader.pages if p.extract_text()])
    elif name.endswith(".txt"):
        text = raw.decode("utf-8", errors="ignore")
    elif name.endswith(".docx"):
        doc = DocxDocument(bio)
        text = "\n".join([p.text for p in doc.paragraphs if p.text])
    elif name.endswith(".pptx"):
        prs = PptxPresentation(bio)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    return splitter.split_text(text)
