# main.py
import os
import requests
from io import BytesIO
import streamlit as st
from dotenv import load_dotenv

# File parsing
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation
from PIL import Image

# LangChain / VectorStore / Embeddings / LLM
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

# -------------------------
# Config
# -------------------------
load_dotenv()
OCR_SPACE_API_KEY = os.getenv("OCR_SPACE_API_KEY")
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY", "")
GROQ_API_KEY = os.getenv("GROQ_API_KEY", "")

st.set_page_config(
    page_title="Gemini + Groq Multi-file Chatbot (FAISS + OCR.Space)",
    page_icon="🤖",
    layout="wide"
)

# -------------------------
# Fungsi cek validitas API key
# -------------------------
def check_google_api_key(key: str) -> bool:
    if not key:
        return False
    try:
        resp = requests.get(
            "https://generativelanguage.googleapis.com/v1/models/gemini-pro",
            params={"key": key},
            timeout=8,
        )
        return resp.status_code == 200
    except:
        return False

def check_groq_api_key(key: str) -> bool:
    if not key:
        return False
    try:
        resp = requests.get(
            "https://api.groq.com/openai/v1/models",
            headers={"Authorization": f"Bearer {key}"},
            timeout=8,
        )
        return resp.status_code == 200
    except:
        return False

# -------------------------
# Validasi API Key dari .env
# -------------------------
valid_google = check_google_api_key(GOOGLE_API_KEY)
valid_groq = check_groq_api_key(GROQ_API_KEY)

# -------------------------
# Sidebar input hanya jika invalid/kosong
# -------------------------
if not valid_google or not valid_groq:
    st.sidebar.header("🔑 API Keys")

    if not valid_google:
        if not GOOGLE_API_KEY:
            st.sidebar.warning("⚠️ GOOGLE_API_KEY belum diisi atau kosong.")
        else:
            st.sidebar.error("❌ GOOGLE_API_KEY tidak valid atau sudah expired. Buat API KEY baru pada https://aistudio.google.com/apikey")
        GOOGLE_API_KEY_INPUT = st.sidebar.text_input(
            "Masukkan GOOGLE_API_KEY (Gemini)", type="password", value=""
        )
        if GOOGLE_API_KEY_INPUT.strip():
            GOOGLE_API_KEY = GOOGLE_API_KEY_INPUT.strip()
            valid_google = check_google_api_key(GOOGLE_API_KEY)
            if valid_google:
                st.sidebar.success("✅ GOOGLE_API_KEY baru valid.")

    if not valid_groq:
        if not GROQ_API_KEY:
            st.sidebar.warning("⚠️ GROQ_API_KEY belum diisi atau kosong.")
        else:
            st.sidebar.error("❌ GROQ_API_KEY tidak valid atau sudah expired. Buat API KEY baru pada https://console.groq.com/keys")
        GROQ_API_KEY_INPUT = st.sidebar.text_input(
            "Masukkan GROQ_API_KEY (Groq)", type="password", value=""
        )
        if GROQ_API_KEY_INPUT.strip():
            GROQ_API_KEY = GROQ_API_KEY_INPUT.strip()
            valid_groq = check_groq_api_key(GROQ_API_KEY)
            if valid_groq:
                st.sidebar.success("✅ GROQ_API_KEY baru valid.")

# -------------------------
# Hentikan jika tidak ada key valid
# -------------------------
if not (valid_google or valid_groq):
    st.error("❌ Tidak ada API Key valid. Tambahkan di .env atau input di sidebar.")
    st.stop()

# -------------------------
# Embeddings
# -------------------------
EMBEDDINGS = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
SPLITTER = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=120)

# -------------------------
# File extractors
# -------------------------
def extract_text_from_pdf(file_bytes: BytesIO):
    text = ""
    try:
        reader = PdfReader(file_bytes)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak PDF: {e}")
    return text

def extract_text_from_txt(file_bytes: BytesIO):
    try:
        return file_bytes.read().decode("utf-8", errors="ignore")
    except Exception as e:
        st.warning(f"⚠️ Gagal baca TXT: {e}")
        return ""

def extract_text_from_docx(file_bytes: BytesIO):
    text = ""
    try:
        file_bytes.seek(0)
        doc = DocxDocument(file_bytes)
        for p in doc.paragraphs:
            if p.text:
                text += p.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak DOCX: {e}")
    return text

def extract_text_from_pptx(file_bytes: BytesIO):
    text = ""
    try:
        file_bytes.seek(0)
        prs = PptxPresentation(file_bytes)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak PPTX: {e}")
    return text

# -------------------------
# OCR.Space Extractor (Image Files)
# -------------------------
def extract_text_from_image(file_bytes: BytesIO, filename="upload.png"):
    if not OCR_SPACE_API_KEY:
        st.warning("⚠️ OCR_SPACE_API_KEY tidak ditemukan di .env")
        return ""

    try:
        file_bytes.seek(0)
        response = requests.post(
            "https://api.ocr.space/parse/image",
            files={"file": (filename, file_bytes, "image/png")},
            data={"apikey": OCR_SPACE_API_KEY, "language": "eng"},
        )
        result = response.json()
        if result.get("IsErroredOnProcessing"):
            st.warning("⚠️ OCR.Space gagal: " + str(result.get("ErrorMessage", ['Unknown error'])))
            return ""
        text = "\n".join([p["ParsedText"] for p in result.get("ParsedResults", []) if "ParsedText" in p])
        return text.strip()
    except Exception as e:
        st.warning(f"⚠️ OCR.Space error: {e}")
        return ""

# -------------------------
# Generic extractor
# -------------------------
def extract_text_from_file(uploaded_file):
    name = uploaded_file.name.lower()
    raw = uploaded_file.read()
    bio = BytesIO(raw)

    if name.endswith(".pdf"):
        return extract_text_from_pdf(bio)
    elif name.endswith(".txt"):
        return extract_text_from_txt(BytesIO(raw))
    elif name.endswith(".docx"):
        return extract_text_from_docx(BytesIO(raw))
    elif name.endswith(".pptx"):
        return extract_text_from_pptx(BytesIO(raw))
    elif name.endswith((".jpg", ".jpeg", ".png", ".gif", ".bmp", ".jfif")):
        return extract_text_from_image(BytesIO(raw), filename=uploaded_file.name)
    elif name.endswith(".doc") or name.endswith(".ppt"):
        st.warning(f"⚠️ File `{uploaded_file.name}` berformat lama (.doc/.ppt). Silakan konversi ke .docx/.pptx.")
        return ""
    else:
        st.warning(f"⚠️ Tipe file `{uploaded_file.name}` tidak didukung.")
        return ""

# -------------------------
# Build documents & FAISS
# -------------------------
def build_documents_from_uploads(uploaded_files):
    docs = []
    for f in uploaded_files:
        text = extract_text_from_file(f)
        if not text or not text.strip():
            continue
        chunks = SPLITTER.split_text(text)
        for i, chunk in enumerate(chunks):
            docs.append(Document(page_content=chunk, metadata={"source_file": f.name, "chunk_id": i}))
    return docs

def build_faiss_from_documents(docs):
    if not docs:
        return None
    vs = FAISS.from_documents(docs, embedding=EMBEDDINGS)
    return vs

# -------------------------
# Prompt formatting helpers
# -------------------------
def format_context(snippets):
    parts = []
    for idx, d in enumerate(snippets, start=1):
        src = d.metadata.get("source_file", "unknown")
        cid = d.metadata.get("chunk_id", "-")
        parts.append(f"[{idx}] ({src}#chunk-{cid})\n{d.page_content}")
    return "\n\n---\n\n".join(parts)

def render_sources(snippets):
    with st.expander("🔎 Sumber konteks yang dipakai"):
        for i, d in enumerate(snippets, start=1):
            src = d.metadata.get("source_file", "unknown")
            cid = d.metadata.get("chunk_id", "-")
            preview = d.page_content[:300].replace("\n", " ")
            st.markdown(f"**[{i}]** **{src}** (chunk {cid})")
            st.caption(preview + ("..." if len(d.page_content) > 300 else ""))

# -------------------------
# Streamlit UI
# -------------------------
st.title("🤖 Gemini 2.5 Flash + Groq Chatbot — Multi-files + OCR.Space")
st.write("Upload banyak file (PDF, TXT, DOCX, PPTX, Images). Gambar akan diproses dengan OCR.Space API.")

# Sidebar upload
st.sidebar.header("📂 Upload & Build")
uploaded_files = st.sidebar.file_uploader(
    "Upload files (pdf, txt, docx, pptx, images) — boleh banyak",
    type=["pdf", "txt", "docx", "pptx", "jpg", "jpeg", "png", "gif", "bmp", "jfif"],
    accept_multiple_files=True
)
build_btn = st.sidebar.button("🚀 Build Vector Store")
clear_btn = st.sidebar.button("🧹 Reset vector store")

if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "indexed_files" not in st.session_state:
    st.session_state.indexed_files = []

if clear_btn:
    st.session_state.vector_store = None
    st.session_state.indexed_files = []
    st.success("Vector store di-reset.")

if build_btn:
    if not uploaded_files:
        st.sidebar.warning("Silakan upload minimal 1 file terlebih dahulu.")
    else:
        with st.spinner("📦 Memproses file dan membuat vector store..."):
            docs = build_documents_from_uploads(uploaded_files)
            if not docs:
                st.sidebar.error("Tidak ada teks valid berhasil diekstrak. Periksa file.")
            else:
                vs = build_faiss_from_documents(docs)
                st.session_state.vector_store = vs
                st.session_state.indexed_files = [f.name for f in uploaded_files]
                st.sidebar.success(f"Vector store terbangun. Dokumen: {len(st.session_state.indexed_files)} | Chunk total: {len(docs)}")

# Show indexed files
if st.session_state.indexed_files:
    st.markdown("**Dokumen terindeks:**")
    st.write(" • " + "\n • ".join(st.session_state.indexed_files))

# -------------------------
# Pilih LLM Provider
# -------------------------
model_choice = st.sidebar.radio(
    "Pilih LLM Provider:",
    ["Gemini 2.5 Flash (Google)", "Groq (llama-3.3-70b-versatile)"]
)

# -------------------------
# Query area
# -------------------------
prompt = st.text_input(
    "Tanyakan sesuatu berdasarkan dokumen yang diupload:",
    placeholder="Misal: Ringkas dokumen tentang topik X..."
)
ask_btn = st.button("Tanyakan")

if ask_btn:
    if not prompt.strip():
        st.warning("Masukkan pertanyaan terlebih dahulu.")
    elif st.session_state.vector_store is None:
        st.info("Belum ada vector store. Upload file dan klik 'Build Vector Store'.")
    else:
        with st.spinner("🔎 Mengambil konteks dari vector store..."):
            results = st.session_state.vector_store.similarity_search(prompt, k=5)

        context_text = format_context(results)
        system_instructions = (
            "Jawablah secara akurat, jelas dan detil sehingga mudah dipahami berdasarkan konteks berikut. "
            "Jika jawaban tidak ada, katakan: "
            "\"Jawaban tidak tersedia dalam konteks yang diberikan, "
            "pelajari dan berikan tambahan informasi yang diperlukan secara detil sesuai konteks atau topik  "
            "yang ditanyakan dari berbagai sumber atau referensi di internet dengan menyertakan link website.\" "
            "Sertakan referensi [angka] ke potongan konteks bila relevan."
        )

        composed_prompt = (
            f"{system_instructions}\n\n"
            f"=== KONTEX ===\n{context_text}\n\n"
            f"=== PERTANYAAN ===\n{prompt}\n\n"
            f"=== JAWABAN ==="
        )

        try:
            if model_choice.startswith("Gemini"):
                from langchain_google_genai import ChatGoogleGenerativeAI
                llm = ChatGoogleGenerativeAI(
                    model="gemini-2.5-flash",
                    temperature=0.2,
                    google_api_key=GOOGLE_API_KEY
                )
                with st.spinner("🤖 Gemini sedang menjawab..."):
                    response = llm.invoke(composed_prompt)
            else:
                from langchain_groq import ChatGroq
                llm = ChatGroq(
                    temperature=0.2,
                    groq_api_key=GROQ_API_KEY,
                    model_name="llama-3.3-70b-versatile"
                )
                with st.spinner("⚡ Groq sedang menjawab..."):
                    response = llm.invoke(composed_prompt)

            st.subheader("💬 Jawaban")
            out_text = getattr(response, "content", None) or str(response)
            st.write(out_text)
            render_sources(results)

        except Exception as e:
            st.error(f"❌ Error saat memanggil LLM: {e}")
