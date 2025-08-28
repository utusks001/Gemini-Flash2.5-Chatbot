# streamlit-app.py
# Gemini Flash 2.5 Chatbot - Upload Multi-files, Build Vector Store and Asking (with Groq Support)
# app.py
import os
from io import BytesIO
import streamlit as st
from dotenv import load_dotenv
from typing import List

# File parsing
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation

# LangChain / VectorStore / Embeddings / LLM
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain_groq import ChatGroq # Import for Groq

# -------------------------
# Config / env
# -------------------------
load_dotenv()
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GROQ_API_KEY = os.getenv("GROQ_API_KEY") # Get Groq API Key

st.set_page_config(
    page_title="Gemini/Groq Multi-file Chatbot (FAISS)", # Updated title
    page_icon="🤖",
    layout="wide"
)

# Check for API keys
if not GOOGLE_API_KEY and not GROQ_API_KEY:
    st.error("❌ GOOGLE_API_KEY atau GROQ_API_KEY tidak ditemukan. Tambahkan salah satu atau keduanya ke file .env sebelum menjalankan.")
    st.stop()

# Embeddings (HuggingFace yang stabil di Streamlit Cloud)
EMBEDDINGS = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

# Text splitter
SPLITTER = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=120)

# -------------------------
# File extractors
# -------------------------
def extract_text_from_pdf(file_bytes: BytesIO) -> str:
    text = ""
    try:
        reader = PdfReader(file_bytes)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak PDF: {e}")
    return text


def extract_text_from_txt(file_bytes: BytesIO) -> str:
    try:
        return file_bytes.read().decode("utf-8", errors="ignore")
    except Exception as e:
        st.warning(f"⚠️ Gagal baca TXT: {e}")
        return ""


def extract_text_from_docx(file_bytes: BytesIO) -> str:
    text = ""
    try:
        file_bytes.seek(0)
        doc = DocxDocument(file_bytes)
        for p in doc.paragraphs:
            if p.text:
                text += p.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak DOCX: {e}")
    return text


def extract_text_from_pptx(file_bytes: BytesIO) -> str:
    text = ""
    try:
        file_bytes.seek(0)
        prs = PptxPresentation(file_bytes)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak PPTX: {e}")
    return text


def extract_text_from_file(uploaded_file) -> str:
    """Best-effort generic extractor"""
    name = uploaded_file.name.lower()
    raw = uploaded_file.read()
    bio = BytesIO(raw)

    if name.endswith(".pdf"):
        return extract_text_from_pdf(bio)
    elif name.endswith(".txt"):
        return extract_text_from_txt(BytesIO(raw))
    elif name.endswith(".docx"):
        return extract_text_from_docx(BytesIO(raw))
    elif name.endswith(".pptx"):
        return extract_text_from_pptx(BytesIO(raw))
    elif name.endswith(".doc") or name.endswith(".ppt"):
        st.warning(f"⚠️ File `{uploaded_file.name}` berformat lama (.doc/.ppt). Silakan konversi ke .docx/.pptx untuk ekstraksi teks yang lebih baik.")
        return ""
    else:
        st.warning(f"⚠️ Tipe file `{uploaded_file.name}` tidak didukung.")
        return ""

# -------------------------
# Build documents & FAISS
# -------------------------
def build_documents_from_uploads(uploaded_files) -> List[Document]:
    docs: List[Document] = []
    for f in uploaded_files:
        text = extract_text_from_file(f)
        if not text or not text.strip():
            continue
        chunks = SPLITTER.split_text(text)
        for i, chunk in enumerate(chunks):
            docs.append(Document(page_content=chunk, metadata={"source_file": f.name, "chunk_id": i}))
    return docs


def build_faiss_from_documents(docs: List[Document]) -> FAISS | None:
    if not docs:
        return None
    vs = FAISS.from_documents(docs, embedding=EMBEDDINGS)
    return vs

# -------------------------
# Prompt formatting helpers
# -------------------------
def format_context(snippets: List[Document]) -> str:
    parts = []
    for idx, d in enumerate(snippets, start=1):
        src = d.metadata.get("source_file", "unknown")
        cid = d.metadata.get("chunk_id", "-")
        parts.append(f"[{idx}] ({src}#chunk-{cid})\n{d.page_content}")
    return "\n\n---\n\n".join(parts)


def render_sources(snippets: List[Document]):
    with st.expander("🔎 Sumber konteks yang dipakai"):
        for i, d in enumerate(snippets, start=1):
            src = d.metadata.get("source_file", "unknown")
            cid = d.metadata.get("chunk_id", "-")
            preview = d.page_content[:300].replace("\n", " ")
            st.markdown(f"**[{i}]** **{src}** (chunk {cid})")
            st.caption(preview + ("..." if len(d.page_content) > 300 else ""))

# -------------------------
# Streamlit UI
# -------------------------
st.title("🤖 Gemini/Groq Chatbot — Multi-files, Build Vector Store and Asking") # Updated title
st.write("Upload banyak file (PDF, TXT, DOCX, PPTX). Untuk .doc/.ppt (format lama), silakan convert ke .docx/.pptx jika ekstraksi kosong.")

# Sidebar
st.sidebar.header("📂 Upload & Build")
uploaded_files = st.sidebar.file_uploader(
    "Upload files (pdf, txt, docx, pptx) — boleh banyak",
    type=["pdf", "txt", "docx", "pptx"],
    accept_multiple_files=True
)
build_btn = st.sidebar.button("🚀 Build Vector Store")
clear_btn = st.sidebar.button("🧹 Reset vector store")

# Model selection in sidebar
st.sidebar.header("⚙️ Pengaturan Model")
model_options = []
if GOOGLE_API_KEY:
    model_options.append("Gemini 2.5 Flash")
if GROQ_API_KEY:
    model_options.extend(["llama3-8b-8192", "llama3-70b-8192", "mixtral-8x7b-32768", "gemma-7b-it"]) # Example Groq models
if not model_options:
    st.sidebar.warning("Tidak ada API Key yang ditemukan. Silakan tambahkan GOOGLE_API_KEY atau GROQ_API_KEY di file .env")
    st.stop()

selected_model = st.sidebar.selectbox("Pilih Model LLM", model_options)


if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "indexed_files" not in st.session_state:
    st.session_state.indexed_files = []

if clear_btn:
    st.session_state.vector_store = None
    st.session_state.indexed_files = []
    st.success("Vector store di-reset.")

if build_btn:
    if not uploaded_files:
        st.sidebar.warning("Silakan upload minimal 1 file terlebih dahulu.")
    else:
        with st.spinner("📦 Memproses file dan membuat vector store..."):
            docs = build_documents_from_uploads(uploaded_files)
            if not docs:
                st.sidebar.error("Tidak ada teks valid yang berhasil diekstrak. Periksa file atau konversi ke docx/pptx.")
            else:
                vs = build_faiss_from_documents(docs)
                st.session_state.vector_store = vs
                st.session_state.indexed_files = [f.name for f in uploaded_files]
                st.sidebar.success(f"Vector store terbangun. Dokumen terindeks: {len(st.session_state.indexed_files)} | Chunk total: {len(docs)}")

# Show indexed files
if st.session_state.indexed_files:
    st.markdown("**Dokumen terindeks:**")
    st.write(" • " + "\n • ".join(st.session_state.indexed_files))

# Query area
prompt = st.text_input(
    "Tanyakan sesuatu berdasarkan dokumen yang diupload:",
    placeholder="Misal: Ringkas dokumen tentang topik X..."
)
ask_btn = st.button("Tanyakan")

if ask_btn:
    if not prompt or not prompt.strip():
        st.warning("Masukkan pertanyaan terlebih dahulu.")
    elif st.session_state.vector_store is None:
        st.info("Belum ada vector store. Upload file dan klik 'Build Vector Store' di sidebar.")
    else:
        with st.spinner("🔎 Mengambil konteks dari vector store..."):
            results = st.session_state.vector_store.similarity_search(prompt, k=5)

        context_text = format_context(results)
        system_instructions = (
            "Jawablah seakurat dan sedetil mungkin sehingga mudah dipahami berdasarkan konteks berikut. "
            "Jika jawaban tidak ada, katakan: "
            "\"Jawaban tidak tersedia dalam konteks yang diberikan, "
            "pelajari dan berikan tambahan informasi yang diperlukan sesuai konteks atau topik  "
            "yang ditanyakan dari berbagai sumber atau referensi di internet dengan menyertakan link website.\" "
            "Sertakan referensi [angka] ke potongan konteks bila relevan."
        )

        composed_prompt = (
            f"{system_instructions}\n\n"
            f"=== KONTEX ===\n{context_text}\n\n"
            f"=== PERTANYANAN ===\n{prompt}\n\n"
            f"=== JAWABAN ==="
        )

        # Initialize LLM based on selection
        llm = None
        if selected_model == "Gemini 2.5 Flash":
            if GOOGLE_API_KEY:
                llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.2)
            else:
                st.error("GOOGLE_API_KEY tidak ditemukan untuk model Gemini.")
                st.stop()
        elif selected_model in ["llama3-8b-8192", "llama3-70b-8192", "mixtral-8x7b-32768", "gemma-7b-it"]:
            if GROQ_API_KEY:
                llm = ChatGroq(temperature=0.2, groq_api_key=GROQ_API_KEY, model_name=selected_model)
            else:
                st.error("GROQ_API_KEY tidak ditemukan untuk model Groq.")
                st.stop()
        else:
            st.error("Model yang dipilih tidak valid atau API Key tidak tersedia.")
            st.stop()

        try:
            with st.spinner(f"🤖 {selected_model} sedang menjawab..."):
                response = llm.invoke(composed_prompt)

            st.subheader("💬 Jawaban")
            out_text = getattr(response, "content", None) or (
                response.candidates[0].content if getattr(response, "candidates", None) else str(response)
            )
            st.write(out_text)
            render_sources(results)

        except Exception as e:
            st.error(f"❌ Error saat memanggil {selected_model}: {e}")
