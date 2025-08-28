# app.py
# app.py
import os
from io import BytesIO
from typing import List, Optional

import streamlit as st
from dotenv import load_dotenv

# File parsing
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation as PptxPresentation

# LangChain / embeddings / vectorstore
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS

# Chat/LLM integrations
# Google Generative (Gemini)
from langchain_google_genai import ChatGoogleGenerativeAI

# Llama (local via llama.cpp python binding)
# Note: llama-cpp-python and langchain_community integration
try:
    from langchain_community.llms import LlamaCpp
    LLAMA_AVAILABLE = True
except Exception:
    LLAMA_AVAILABLE = False

# Groq (direct SDK)
GROQ_AVAILABLE = False
try:
    # groq official SDK (simple direct usage)
    from groq.cloud.core import ChatCompletion as GroqChatCompletion
    GROQ_AVAILABLE = True
except Exception:
    GROQ_AVAILABLE = False

# Load .env
load_dotenv()

# Config
st.set_page_config(page_title="Gemini+Groq+Llama Multi-file Chatbot", page_icon="🤖", layout="wide")
st.title("🤖 Multi-backend Multi-file Chatbot — Gemini / Groq / Llama")
st.markdown(
    "Upload file (PDF / TXT / DOCX / PPTX). Build FAISS vector store dari dokumen, lalu pilih backend LLM untuk menjawab "
    "(Google Gemini, Groq Cloud, atau Llama lokal via llama.cpp)."
)

# Environment variables (optional)
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
# Accept both names depending on SDK/version
GROQ_API_KEY = os.getenv("GROQ_API_KEY") or os.getenv("GROQ_SECRET_ACCESS_KEY") or os.getenv("GROQ_KEY")
LLAMA_MODEL_PATH = os.getenv("LLAMA_MODEL_PATH")  # path to gguf / ggml model file if using llama-cpp-python

# Embeddings & splitter
EMBEDDINGS = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")
SPLITTER = RecursiveCharacterTextSplitter(chunk_size=800, chunk_overlap=120)

# -------------------------
# File extractors
# -------------------------
def extract_text_from_pdf(file_bytes: BytesIO) -> str:
    text = ""
    try:
        reader = PdfReader(file_bytes)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak PDF: {e}")
    return text

def extract_text_from_txt(file_bytes: BytesIO) -> str:
    try:
        return file_bytes.read().decode("utf-8", errors="ignore")
    except Exception as e:
        st.warning(f"⚠️ Gagal baca TXT: {e}")
        return ""

def extract_text_from_docx(file_bytes: BytesIO) -> str:
    text = ""
    try:
        file_bytes.seek(0)
        doc = DocxDocument(file_bytes)
        for p in doc.paragraphs:
            if p.text:
                text += p.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak DOCX: {e}")
    return text

def extract_text_from_pptx(file_bytes: BytesIO) -> str:
    text = ""
    try:
        file_bytes.seek(0)
        prs = PptxPresentation(file_bytes)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception as e:
        st.warning(f"⚠️ Gagal ekstrak PPTX: {e}")
    return text

def extract_text_from_file(uploaded_file) -> str:
    name = uploaded_file.name.lower()
    raw = uploaded_file.read()
    bio = BytesIO(raw)

    if name.endswith(".pdf"):
        return extract_text_from_pdf(bio)
    elif name.endswith(".txt"):
        return extract_text_from_txt(BytesIO(raw))
    elif name.endswith(".docx"):
        return extract_text_from_docx(BytesIO(raw))
    elif name.endswith(".pptx"):
        return extract_text_from_pptx(BytesIO(raw))
    elif name.endswith(".doc") or name.endswith(".ppt"):
        st.warning(f"⚠️ File `{uploaded_file.name}` berformat lama (.doc/.ppt). Silakan konversi ke .docx/.pptx.")
        return ""
    else:
        st.warning(f"⚠️ Tipe file `{uploaded_file.name}` tidak didukung.")
        return ""

# -------------------------
# Build documents & FAISS
# -------------------------
def build_documents_from_uploads(uploaded_files) -> List[Document]:
    docs: List[Document] = []
    for f in uploaded_files:
        text = extract_text_from_file(f)
        if not text or not text.strip():
            continue
        chunks = SPLITTER.split_text(text)
        for i, chunk in enumerate(chunks):
            docs.append(Document(page_content=chunk, metadata={"source_file": f.name, "chunk_id": i}))
    return docs

def build_faiss_from_documents(docs: List[Document]) -> Optional[FAISS]:
    if not docs:
        return None
    vs = FAISS.from_documents(docs, embedding=EMBEDDINGS)
    return vs

# -------------------------
# Helpers: format context / render sources
# -------------------------
def format_context(snippets: List[Document]) -> str:
    parts = []
    for idx, d in enumerate(snippets, start=1):
        src = d.metadata.get("source_file", "unknown")
        cid = d.metadata.get("chunk_id", "-")
        parts.append(f"[{idx}] ({src}#chunk-{cid})\n{d.page_content}")
    return "\n\n---\n\n".join(parts)

def render_sources(snippets: List[Document]):
    with st.expander("🔎 Sumber konteks yang dipakai"):
        for i, d in enumerate(snippets, start=1):
            src = d.metadata.get("source_file", "unknown")
            cid = d.metadata.get("chunk_id", "-")
            preview = d.page_content[:300].replace("\n", " ")
            st.markdown(f"**[{i}]** **{src}** (chunk {cid})")
            st.caption(preview + ("..." if len(d.page_content) > 300 else ""))

# -------------------------
# UI: sidebar
# -------------------------
st.sidebar.header("📂 Upload & Build")
uploaded_files = st.sidebar.file_uploader(
    "Upload files (pdf, txt, docx, pptx) — boleh banyak",
    type=["pdf", "txt", "docx", "pptx"],
    accept_multiple_files=True
)
build_btn = st.sidebar.button("🚀 Build Vector Store")
clear_btn = st.sidebar.button("🧹 Reset vector store")

# backend choices
st.sidebar.markdown("### Pilih backend LLM")
backend = st.sidebar.radio("Backend LLM", options=["google_gemini", "groq", "llama_local"])

# Llama model path input (overrides env if provided)
if backend == "llama_local":
    llama_path_input = st.sidebar.text_input("Path ke model Llama (gguf/ggml)", value=LLAMA_MODEL_PATH or "")

# show detected availability
st.sidebar.markdown("**Ketersediaan SDK**")
st.sidebar.write(f"- Groq SDK: {'✅' if GROQ_AVAILABLE else '❌ (tidak terinstal)'}")
st.sidebar.write(f"- Llama (llama-cpp-python): {'✅' if LLAMA_AVAILABLE else '❌ (tidak terinstal)'}")
st.sidebar.write(f"- Google API key present: {'✅' if GOOGLE_API_KEY else '❌'}")

# session state
if "vector_store" not in st.session_state:
    st.session_state.vector_store = None
if "indexed_files" not in st.session_state:
    st.session_state.indexed_files = []

if clear_btn:
    st.session_state.vector_store = None
    st.session_state.indexed_files = []
    st.success("Vector store di-reset.")

if build_btn:
    if not uploaded_files:
        st.sidebar.warning("Silakan upload minimal 1 file terlebih dahulu.")
    else:
        with st.spinner("📦 Memproses file dan membuat vector store..."):
            docs = build_documents_from_uploads(uploaded_files)
            if not docs:
                st.sidebar.error("Tidak ada teks valid yang berhasil diekstrak.")
            else:
                vs = build_faiss_from_documents(docs)
                st.session_state.vector_store = vs
                st.session_state.indexed_files = [f.name for f in uploaded_files]
                st.sidebar.success(f"Vector store terbangun. Dokumen terindeks: {len(st.session_state.indexed_files)} | Chunk total: {len(docs)}")

# show indexed files
if st.session_state.indexed_files:
    st.markdown("**Dokumen terindeks:**")
    st.write(" • " + "\n • ".join(st.session_state.indexed_files))

# -------------------------
# Query
# -------------------------
prompt = st.text_input("Tanyakan sesuatu berdasarkan dokumen yang diupload:", placeholder="Misal: Ringkas dokumen tentang topik X...")
ask_btn = st.button("Tanyakan")

if ask_btn:
    if not prompt or not prompt.strip():
        st.warning("Masukkan pertanyaan terlebih dahulu.")
    elif st.session_state.vector_store is None:
        st.info("Belum ada vector store. Upload file dan klik 'Build Vector Store' di sidebar.")
    else:
        with st.spinner("🔎 Mengambil konteks dari vector store..."):
            results = st.session_state.vector_store.similarity_search(prompt, k=5)

        context_text = format_context(results)
        system_instructions = (
            "Jawablah seakurat dan sedetil mungkin berdasarkan konteks berikut. "
            "Jika konteks tidak cukup, jelaskan keterbatasan dan berikan rekomendasi sumber tambahan."
        )

        composed_prompt = (
            f"{system_instructions}\n\n"
            f"=== KONTEX ===\n{context_text}\n\n"
            f"=== PERTANYAAN ===\n{prompt}\n\n"
            f"=== JAWABAN ==="
        )

        # -------- Choose backend and call appropriate model ----------
        try:
            if backend == "google_gemini":
                if not GOOGLE_API_KEY:
                    st.error("❌ GOOGLE_API_KEY tidak ditemukan. Set environment variable atau isi di .env.")
                else:
                    llm = ChatGoogleGenerativeAI(model="gemini-2.5-flash", temperature=0.2, google_api_key=GOOGLE_API_KEY)
                    with st.spinner("🤖 Menghubungi Google Gemini..."):
                        response = llm.invoke(composed_prompt)
                    out_text = getattr(response, "content", None) or (
                        response.candidates[0].content if getattr(response, "candidates", None) else str(response)
                    )
                    st.subheader("💬 Jawaban (Google Gemini)")
                    st.write(out_text)
                    render_sources(results)

            elif backend == "groq":
                if not GROQ_AVAILABLE:
                    st.error("❌ Groq SDK tidak terinstal. Install `groq` atau `langchain-groq` dan restart.")
                elif not GROQ_API_KEY:
                    st.error("❌ GROQ API key tidak ditemukan. Set environment variable GROQ_API_KEY or GROQ_SECRET_ACCESS_KEY.")
                else:
                    # direct minimal usage of groq Cloud SDK
                    with st.spinner("🤖 Menghubungi Groq Cloud..."):
                        # choose a model name you have access to e.g. "llama2-13b" or provider model name
                        model_name = st.sidebar.text_input("Groq model (contoh: llama2-13b-2048)", value="llama2-13b-2048")
                        # Use ChatCompletion context if available
                        chat = GroqChatCompletion(model_name, api_key=GROQ_API_KEY)
                        # send_chat returns (response_text, _, _)
                        resp, _, _ = chat.send_chat(composed_prompt)
                        st.subheader("💬 Jawaban (Groq)")
                        st.write(resp)
                        render_sources(results)

            elif backend == "llama_local":
                model_path = llama_path_input or LLAMA_MODEL_PATH
                if not LLAMA_AVAILABLE:
                    st.error("❌ llama-cpp-python (langchain integration) tidak terinstal. Install `llama-cpp-python`.")
                elif not model_path:
                    st.error("❌ Path model Llama tidak diberikan. Set LLAMA_MODEL_PATH di .env atau isi field di sidebar.")
                else:
                    with st.spinner("🤖 Menjalankan Llama lokal (llama.cpp)..."):
                        # instantiate LlamaCpp (langchain_community)
                        llm = LlamaCpp(model_path=model_path, n_ctx=2048, temperature=0.2)
                        answer = llm(composed_prompt)
                        st.subheader("💬 Jawaban (Llama lokal)")
                        st.write(answer)
                        render_sources(results)
            else:
                st.error("Backend tidak dikenal.")
        except Exception as e:
            st.error(f"❌ Error saat memanggil backend LLM: {e}")
